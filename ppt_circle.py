import builtins
import collections.abc as collections
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

# Create a new presentation
prs = Presentation()

# Add a blank slide
slide_layout = prs.slide_layouts[6]  # 6 is the index for a blank slide
slide = prs.slides.add_slide(slide_layout)

# Add a circle shape
left = top = Inches(1.0)
width = height = Inches(1.0)
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)

# Save the presentation
prs.save('circle.pptx')
