import collections.abc as collections
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Add a blank slide
slide_layout = prs.slide_layouts[6]  # 6 is the index for a blank slide
slide = prs.slides.add_slide(slide_layout)

# Add a blue rectangle for the sky
sky = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.5))
sky.fill.solid()
sky.fill.fore_color.rgb = RGBColor(135, 206, 235)  # Sky blue color
sky.line.fill.background()

# Add a yellow circle for the sun
sun = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7), Inches(0.5), Inches(1), Inches(1))
sun.fill.solid()
sun.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow color
sun.line.fill.background()

# Add a green rectangle for the ground
ground = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.5), Inches(10), Inches(1.5))
ground.fill.solid()
ground.fill.fore_color.rgb = RGBColor(34, 139, 34)  # Green color
ground.line.fill.background()

# Add white clouds
for left, top in [(Inches(1), Inches(1)), (Inches(4), Inches(0.5)), (Inches(6), Inches(2))]:
    cloud = slide.shapes.add_shape(MSO_SHAPE.CLOUD, left, top, Inches(1.5), Inches(1))
    cloud.fill.solid()
    cloud.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White color
    cloud.line.fill.background()

# Add flowers
for left in [Inches(1), Inches(3), Inches(5), Inches(7), Inches(9)]:
    flower = slide.shapes.add_shape(MSO_SHAPE.SUN, left, Inches(6), Inches(0.5), Inches(0.5))
    flower.fill.solid()
    flower.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red color
    flower.line.fill.background()

# Add trees
for left in [Inches(2), Inches(5), Inches(8)]:
    trunk = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(5), Inches(0.2), Inches(1))
    trunk.fill.solid()
    trunk.fill.fore_color.rgb = RGBColor(139, 69, 19)  # Brown color
    trunk.line.fill.background()
    leaves = slide.shapes.add_shape(MSO_SHAPE.OVAL, left - Inches(0.4), Inches(4), Inches(1), Inches(1))
    leaves.fill.solid()
    leaves.fill.fore_color.rgb = RGBColor(34, 139, 34)  # Green color
    leaves.line.fill.background()

# Add birds
for left, top in [(Inches(2), Inches(2)), (Inches(3), Inches(1.5)), (Inches(4), Inches(2.5))]:
    bird = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, Inches(0.5), Inches(0.5))
    bird.fill.solid()
    bird.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black color
    bird.line.fill.background()

# Save the presentation
prs.save('sunny_day.pptx')

